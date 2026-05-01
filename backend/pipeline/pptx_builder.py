"""
pptx_builder.py
Build a PowerPoint file from vectorised shape data.

Shapes are drawn as closed cubic-Bezier freeforms using Catmull-Rom
interpolation so that all curves are smooth regardless of point count.

Pipeline per contour
────────────────────
1. Convert pixel coords → absolute EMU on the slide.
2. Compute Catmull-Rom control points for every edge.
3. Find the true bounding box (endpoints + control points can overshoot).
4. Subtract the bbox origin → relative coords for the path.
5. Emit  <a:moveTo> + N×<a:cubicBezTo> + <a:close>  inside <a:path>.
6. Wrap in a full <p:sp> element and append to the slide's shape tree.
"""

import io
from typing import List, Tuple

import numpy as np
from lxml import etree
from pptx import Presentation
from pptx.util import Emu

# ── Slide dimensions (16:9 widescreen) ────────────────────────────────────────
SLIDE_W = 12_192_000  # EMU  (13.33 inches)
SLIDE_H = 6_858_000   # EMU  (7.5 inches)

# ── OOXML namespace URIs ───────────────────────────────────────────────────────
_PML = "http://schemas.openxmlformats.org/presentationml/2006/main"
_DML = "http://schemas.openxmlformats.org/drawingml/2006/main"
_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

_NS = {"p": _PML, "a": _DML, "r": _REL}


# ── Catmull-Rom → cubic Bezier ─────────────────────────────────────────────────

def _catmull_rom_segments(pts: np.ndarray):
    """
    For a closed polygon of N points yield (cp1, cp2, end) for each edge.

    Uniform Catmull-Rom → cubic Bezier conversion:
        CP1 = P[i]   + (P[i+1] - P[i-1]) / 6
        CP2 = P[i+1] - (P[i+2] - P[i])   / 6
    The curve passes through every P[i] and is C1-continuous.
    """
    n = len(pts)
    for i in range(n):
        p0 = pts[(i - 1) % n]
        p1 = pts[i]
        p2 = pts[(i + 1) % n]
        p3 = pts[(i + 2) % n]
        cp1 = p1 + (p2 - p0) / 6.0
        cp2 = p2 - (p3 - p1) / 6.0
        yield cp1, cp2, p2


# ── XML helpers ────────────────────────────────────────────────────────────────

def _pt(x: float, y: float) -> str:
    return f'<a:pt x="{round(x)}" y="{round(y)}"/>'


def _build_sp_xml(
    shape_id: int,
    emu_pts: np.ndarray,
    fill_rgb: Tuple[int, int, int],
) -> bytes:
    """
    Construct the complete <p:sp> XML string for one closed Bezier shape.

    emu_pts  – absolute EMU coordinates (N×2 float array)
    fill_rgb – (R, G, B) 0-255
    """
    n = len(emu_pts)
    r, g, b = fill_rgb
    color_hex = f"{r:02X}{g:02X}{b:02X}"

    # ── Compute Catmull-Rom segments in absolute coords ──────────────────
    segs = list(_catmull_rom_segments(emu_pts))

    # ── True bounding box: endpoints + control points ────────────────────
    # Control points can overshoot the endpoint bbox; include them so the
    # shape's declared size encompasses the full visual area.
    ctrl_flat = np.array([[cp1, cp2] for cp1, cp2, _ in segs]).reshape(-1, 2)
    all_pts = np.vstack([emu_pts, ctrl_flat])
    min_xy = all_pts.min(axis=0)
    max_xy = all_pts.max(axis=0)

    off_x = int(min_xy[0])
    off_y = int(min_xy[1])
    cx = max(1, int(max_xy[0] - min_xy[0]))
    cy = max(1, int(max_xy[1] - min_xy[1]))

    # ── Relative coordinates (local to shape) ────────────────────────────
    rel_pts = emu_pts - min_xy

    # ── Build the <a:path> inner XML ─────────────────────────────────────
    x0, y0 = rel_pts[0]
    path_parts = [f"<a:moveTo>{_pt(x0, y0)}</a:moveTo>"]

    for cp1, cp2, end in _catmull_rom_segments(rel_pts):
        path_parts.append(
            f"<a:cubicBezTo>"
            f"{_pt(*cp1)}{_pt(*cp2)}{_pt(*end)}"
            f"</a:cubicBezTo>"
        )
    path_parts.append("<a:close/>")

    path_xml = (
        f'<a:path w="{cx}" h="{cy}">'
        + "".join(path_parts)
        + "</a:path>"
    )

    # ── Full <p:sp> XML ───────────────────────────────────────────────────
    # Mirrors the structure python-pptx emits for a freeform shape, with
    # solidFill and noFill line added directly in spPr.
    sp_xml = (
        f'<p:sp xmlns:p="{_PML}" xmlns:a="{_DML}" xmlns:r="{_REL}">'
        f"<p:nvSpPr>"
        f'<p:cNvPr id="{shape_id}" name="Shape {shape_id}"/>'
        f"<p:cNvSpPr/>"
        f"<p:nvPr/>"
        f"</p:nvSpPr>"
        f"<p:spPr>"
        f"<a:xfrm>"
        f'<a:off x="{off_x}" y="{off_y}"/>'
        f'<a:ext cx="{cx}" cy="{cy}"/>'
        f"</a:xfrm>"
        f"<a:custGeom>"
        f"<a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>"
        f'<a:rect l="l" t="t" r="r" b="b"/>'
        f"<a:pathLst>{path_xml}</a:pathLst>"
        f"</a:custGeom>"
        f"<a:solidFill><a:srgbClr val=\"{color_hex}\"/></a:solidFill>"
        f"<a:ln><a:noFill/></a:ln>"
        f"</p:spPr>"
        # Style block matches python-pptx's default; spPr fill overrides it.
        f"<p:style>"
        f'<a:lnRef idx="1"><a:schemeClr val="accent1"/></a:lnRef>'
        f'<a:fillRef idx="3"><a:schemeClr val="accent1"/></a:fillRef>'
        f'<a:effectRef idx="2"><a:schemeClr val="accent1"/></a:effectRef>'
        f'<a:fontRef idx="minor"><a:schemeClr val="lt1"/></a:fontRef>'
        f"</p:style>"
        f"<p:txBody>"
        f'<a:bodyPr rtlCol="0" anchor="ctr"/>'
        f"<a:lstStyle/>"
        f'<a:p><a:pPr algn="ctr"/></a:p>'
        f"</p:txBody>"
        f"</p:sp>"
    )
    return sp_xml.encode("utf-8")


# ── Public API ─────────────────────────────────────────────────────────────────

def build_pptx(shapes_data, img_width: int, img_height: int) -> bytes:
    """
    Create a Presentation with one slide populated by smooth Bezier shapes.

    shapes_data – from vectorize.extract_shapes()
    img_width/height – pixel dimensions of the source image
    """
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    sp_tree = slide.shapes._spTree

    # Scale image to fit slide, centred, aspect-ratio preserved
    scale = min(SLIDE_W / img_width, SLIDE_H / img_height)
    off_x_slide = (SLIDE_W - img_width * scale) / 2.0
    off_y_slide = (SLIDE_H - img_height * scale) / 2.0

    def to_emu(px: float, py: float) -> Tuple[float, float]:
        return px * scale + off_x_slide, py * scale + off_y_slide

    # Start IDs above any existing elements in the blank slide
    shape_id = len(slide.shapes) + 10

    for color_rgb, contour_list in shapes_data:
        for points in contour_list:
            if len(points) < 3:
                continue

            emu_pts = np.array([to_emu(px, py) for px, py in points], dtype=float)

            try:
                sp_elem = etree.fromstring(_build_sp_xml(shape_id, emu_pts, color_rgb))
                sp_tree.append(sp_elem)
                shape_id += 1
            except Exception:
                continue

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
